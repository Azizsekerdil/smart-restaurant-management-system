"""Tanıtım sunumunu HTML, PowerPoint ve PDF olarak üretir.

İçerik ``scripts/presentation_content.py`` dosyasındadır; burada yalnızca
biçimlendirme yapılır. Bir cümle değiştiğinde tek yeri düzeltip bu
script'i çalıştırmak yeterlidir.

Üretilen dosyalar (``sunum/`` klasörü)
--------------------------------------
    Akilli_Restaurant_Tanitim.html          ekran sunumu (TR)
    Akilli_Restaurant_Tanitim.pdf
    Akilli_Restaurant_Tanitim.pptx
    Akilli_Restaurant_Tanitim_Baski.pdf     baskı için açık zemin
    Akilli_Restaurant_Tanitim_Baski.pptx
    Akilli_Restaurant_Intro_EN.html         ve İngilizce karşılıkları
    ...

Neden iki sürüm
---------------
Ekran sunumu koyu zeminlidir; projeksiyon ve monitörde okunaklıdır.
Baskı sürümü açık zeminlidir: koyu zeminli slaytları yazdırmak hem toner
harcar hem de metni okunaksız kılar.

Kullanım
--------
    python scripts/make_presentation.py
    python scripts/make_presentation.py --only html
"""

from __future__ import annotations

import html as html_module
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import landscape
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas as pdf_canvas

sys.path.insert(0, str(Path(__file__).resolve().parent))
from presentation_content import BRAND, META, SLIDES, resolve_metrics  # noqa: E402

PROJECT_ROOT = Path(__file__).resolve().parent.parent
#: Yayımlanan sunumun çıktı klasörü. Özgün (dahili) sunum dosyalarının
#: ÜZERİNE YAZILMAZ: public sürüm ayrı bir klasöre, ayrı adla üretilir.
OUTPUT_DIR = PROJECT_ROOT / "docs" / "presentation"

#: 16:9
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
PAGE_SIZE = landscape((SLIDE_W_IN * 72, SLIDE_H_IN * 72))

FILENAMES = {
    "tr": {
        "screen": "Akilli_Restaurant_Tanitim_PUBLIC",
        "print": "Akilli_Restaurant_Tanitim_Baski_PUBLIC",
    },
    "en": {
        "screen": "Smart_Restaurant_Intro_EN_PUBLIC",
        "print": "Smart_Restaurant_Intro_EN_Print_PUBLIC",
    },
}


def use_language(language: str) -> None:
    """Slaytlardaki ``{{ölçüm}}`` yer tutucularını o dil için doldurur.

    Üretici fonksiyonlar modül düzeyindeki ``SLIDES`` listesini okur; bu
    yüzden dil değiştiğinde liste yeniden çözümlenip yerine konur. Böylece
    sunumdaki her sayı, üretim anında depodan ÖLÇÜLMÜŞ olur.
    """
    globals()["SLIDES"] = resolve_metrics(SLIDES_SOURCE, language)


#: Çözümlenmemiş özgün içerik; her dil geçişinde buradan yeniden üretilir.
SLIDES_SOURCE = SLIDES

#: Ekran görüntülerinin kaynağı (``scripts/capture_screenshots.py`` yazar).
SCREENSHOT_DIR = PROJECT_ROOT / "sunum" / "screenshots"


# ==================================================================
#  Renk paletleri
# ==================================================================
def palette(printable: bool) -> dict:
    """Ekran (koyu) ve baskı (açık) paletleri."""
    if printable:
        return {
            "bg": "#FFFFFF",
            "panel": "#F4F6FA",
            "title": "#0B1220",
            "text": "#26303F",
            "muted": "#5B6B80",
            "brand": "#14508C",
            "accent": "#B4530A",
            "line": "#D3DAE4",
            "chip_bg": "#E8EEF7",
        }
    return {
        "bg": BRAND["dark"],
        "panel": BRAND["surface"],
        "title": "#FFFFFF",
        "text": BRAND["text_light"],
        "muted": BRAND["muted"],
        "brand": BRAND["primary"],
        "accent": BRAND["accent"],
        "line": "#22304A",
        "chip_bg": "#17263E",
    }


# ==================================================================
#  Yazı tipi (PDF)
# ==================================================================
#: Seçilen yazı tipinde bulunmayan işaretlerin güvenli karşılıkları.
#: Eksik bir glif PDF'te sessizce boşluk bırakır; bu sözlük onu görünür
#: ve okunur bir karşılığa çevirir.
GLYPH_FALLBACKS = {
    "\u20ba": "TL",  # lira işareti
    "\u2194": "<->",
    "\u2192": "->",
    "\u2713": "+",
}

#: Üretimde kullanılan yazı tipinin adı ve lisansı (THIRD_PARTY_NOTICES.md).
SELECTED_FONT: dict[str, str] = {}


def register_pdf_fonts() -> tuple[str, str]:
    """Türkçe karakterleri gösteren, **özgür lisanslı** bir yazı tipi kaydeder.

    ReportLab'ın yerleşik Helvetica'sı WinAnsi kodlamasını kullanır ve
    ``ş ğ ı İ`` gibi Türkçe harfleri içermez; bu harfler çıktıda kaybolur.

    Yayımlanan bir PDF'e tescilli yazı tipi gömmek lisans sorunu
    yaratabileceğinden sıra özgür lisanslı yazı tiplerindedir:

    1. **DejaVu Sans** (Bitstream Vera türevi, izin verici lisans) — varsa.
    2. **Bitstream Vera Sans** — ReportLab ile birlikte gelir, bu yüzden
       her kurulumda bulunur; ek indirme gerektirmez.
    3. Son çare olarak sistem yazı tipi; bu durumda çıktı dağıtılmadan
       önce lisans doğrulaması gerekir ve uyarı yazılır.
    """
    fonts_dir = Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts"

    # 1) DejaVu Sans (sistemde varsa)
    dejavu_regular, dejavu_bold = fonts_dir / "DejaVuSans.ttf", fonts_dir / "DejaVuSans-Bold.ttf"
    if dejavu_regular.is_file() and dejavu_bold.is_file():
        pdfmetrics.registerFont(TTFont("DejaVuSans", str(dejavu_regular)))
        pdfmetrics.registerFont(TTFont("DejaVuSans-Bold", str(dejavu_bold)))
        SELECTED_FONT.update(name="DejaVu Sans", licence="Bitstream Vera / DejaVu (izin verici)")
        return "DejaVuSans", "DejaVuSans-Bold"

    # 2) Bitstream Vera Sans — ReportLab paketinin içinden.
    try:
        import reportlab

        bundled = Path(reportlab.__file__).resolve().parent / "fonts"
        vera_regular, vera_bold = bundled / "Vera.ttf", bundled / "VeraBd.ttf"
        if vera_regular.is_file() and vera_bold.is_file():
            pdfmetrics.registerFont(TTFont("BitstreamVeraSans", str(vera_regular)))
            pdfmetrics.registerFont(TTFont("BitstreamVeraSans-Bold", str(vera_bold)))
            SELECTED_FONT.update(
                name="Bitstream Vera Sans", licence="Bitstream Vera Fonts Copyright"
            )
            return "BitstreamVeraSans", "BitstreamVeraSans-Bold"
    except Exception as exc:  # noqa: BLE001 - yazı tipi yoksa akış sürsün
        print(f"  [!] Paketli Vera yazı tipi yüklenemedi: {exc}")

    # 3) Sistem yazı tipi (tescilli olabilir).
    for name, regular, bold in [
        ("Arial", "arial.ttf", "arialbd.ttf"),
        ("SegoeUI", "segoeui.ttf", "segoeuib.ttf"),
        ("Calibri", "calibri.ttf", "calibrib.ttf"),
    ]:
        regular_path, bold_path = fonts_dir / regular, fonts_dir / bold
        if regular_path.is_file() and bold_path.is_file():
            pdfmetrics.registerFont(TTFont(name, str(regular_path)))
            pdfmetrics.registerFont(TTFont(f"{name}-Bold", str(bold_path)))
            SELECTED_FONT.update(name=name, licence="TESCILLI — dağıtımdan önce doğrulayın")
            print(
                f"  [!] Özgür lisanslı yazı tipi bulunamadı; '{name}' kullanılıyor. "
                "PDF dağıtılmadan önce yazı tipi lisansını doğrulayın."
            )
            return name, f"{name}-Bold"

    print("  [!] Türkçe destekli yazı tipi bulunamadı; Helvetica'ya düşülüyor.")
    SELECTED_FONT.update(name="Helvetica", licence="ReportLab yerleşik")
    return "Helvetica", "Helvetica-Bold"


# ==================================================================
#  Yardımcılar
# ==================================================================
def safe_glyphs(text: str) -> str:
    """Yazı tipinde bulunmayan işaretleri okunur karşılıklarıyla değiştirir."""
    for source, target in GLYPH_FALLBACKS.items():
        text = text.replace(source, target)
    return text


def strip_markup(text: str) -> str:
    """İçerikteki **vurgu** işaretlerini kaldırır."""
    return safe_glyphs(text.replace("**", ""))


def upper(text: str, language: str) -> str:
    """Dile duyarlı büyük harf.

    Python'un ``str.upper()`` çağrısı Türkçe'de yanlış sonuç verir:
    ``i`` harfini ``I`` yapar, oysa Türkçe'de karşılığı ``İ``'dir
    ("önemli" -> "ÖNEMLI" yerine "ÖNEMLİ" olmalıdır). Tarayıcılar
    ``lang="tr"`` verildiğinde bunu doğru yapar; PDF ve PowerPoint
    çıktılarında elle düzeltmek gerekir.
    """
    if language == "tr":
        return text.replace("i", "İ").replace("ı", "I").upper()
    return text.upper()


def emphasis_parts(text: str) -> list[tuple[str, bool]]:
    """**vurgulu** metni (parça, kalın_mı) listesine ayırır."""
    parts: list[tuple[str, bool]] = []
    for index, chunk in enumerate(re.split(r"\*\*", text)):
        if chunk:
            parts.append((chunk, index % 2 == 1))
    return parts


def wrap_text(text: str, font: str, size: float, max_width: float) -> list[str]:
    """Metni verilen genişliğe sığacak satırlara böler (PDF için)."""
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if pdfmetrics.stringWidth(candidate, font, size) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


# ==================================================================
#  HTML
# ==================================================================
HTML_TEMPLATE = """<!doctype html>
<html lang="{lang}">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
  :root {{
    --bg: {bg}; --panel: {panel}; --title: {title_color}; --text: {text};
    --muted: {muted}; --brand: {brand}; --accent: {accent}; --line: {line};
    --chip: {chip_bg};
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    background: #05080F; color: var(--text); min-height: 100vh;
    font-family: "Segoe UI", system-ui, -apple-system, "Helvetica Neue", Arial, sans-serif;
    display: flex; flex-direction: column; align-items: center; padding: 24px 16px 96px;
  }}
  .deck {{ width: min(1180px, 100%); display: flex; flex-direction: column; gap: 28px; }}
  .slide {{
    background: var(--bg); border: 1px solid var(--line); border-radius: 18px;
    /* 16:9 tercih edilir ama zorunlu değildir: dar ekranda metin sarınca
       slayt uzar. "overflow: hidden" olsaydı içerik sessizce kırpılırdı —
       bir sunumda en kötü hata, cümlenin yarısının kaybolmasıdır. */
    aspect-ratio: 16 / 9; min-height: 0; height: auto; overflow: visible;
    padding: 54px 62px; position: relative;
    display: flex; flex-direction: column; scroll-margin-top: 20px;
  }}
  .slide::after {{
    content: ""; position: absolute; inset: auto 0 0 0; height: 4px;
    background: linear-gradient(90deg, var(--brand), var(--accent));
  }}
  /* Altbilgi mutlak konumlu değil: slayt uzadığında içeriğin üstüne
     binmemesi için akışın sonunda durur. */
  .foot {{
    margin-top: auto; padding-top: 26px; display: flex; justify-content: space-between;
    font-size: 12px; color: var(--muted); letter-spacing: .04em;
  }}
  h1 {{ font-size: 60px; font-weight: 800; color: var(--title); line-height: 1.05; letter-spacing: -.02em; }}
  h2 {{ font-size: 38px; font-weight: 700; color: var(--title); margin-bottom: 6px; letter-spacing: -.01em; }}
  h3 {{ font-size: 17px; font-weight: 700; color: var(--title); }}
  .lead {{ font-size: 19px; color: var(--muted); margin-bottom: 26px; max-width: 70ch; line-height: 1.5; }}
  ul {{ list-style: none; display: flex; flex-direction: column; gap: 14px; }}
  li {{ font-size: 19px; line-height: 1.45; padding-left: 30px; position: relative; }}
  li::before {{
    content: ""; position: absolute; left: 4px; top: .55em; width: 9px; height: 9px;
    border-radius: 3px; background: var(--brand);
  }}
  .cover {{ justify-content: center; gap: 18px; }}
  .cover .mark {{
    width: 76px; height: 76px; border-radius: 22px; display: grid; place-items: center;
    background: linear-gradient(135deg, var(--brand), var(--accent));
    font-size: 38px; margin-bottom: 10px;
  }}
  .cover .sub {{ font-size: 26px; color: var(--muted); font-weight: 500; }}
  .cover .tag {{ font-size: 20px; color: var(--text); max-width: 62ch; line-height: 1.5; margin-top: 8px; }}
  .chips {{ display: flex; gap: 10px; flex-wrap: wrap; margin-top: 26px; }}
  .chip {{
    background: var(--chip); border: 1px solid var(--line); color: var(--text);
    padding: 8px 16px; border-radius: 999px; font-size: 14px; font-weight: 600;
  }}
  .stats {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 18px; margin-top: 12px; }}
  .stat {{ background: var(--panel); border: 1px solid var(--line); border-radius: 14px; padding: 26px 22px; }}
  .stat .v {{ font-size: 52px; font-weight: 800; color: var(--brand); line-height: 1; }}
  .stat .l {{ font-size: 17px; font-weight: 600; color: var(--title); margin-top: 10px; }}
  .stat .n {{ font-size: 13px; color: var(--muted); margin-top: 6px; line-height: 1.4; }}
  .grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 14px; margin-top: 6px; }}
  .card {{ background: var(--panel); border: 1px solid var(--line); border-radius: 12px; padding: 16px 16px 18px; }}
  .card h3 {{ font-size: 15px; margin-bottom: 5px; }}
  .card p {{ font-size: 13px; color: var(--muted); line-height: 1.45; }}
  .split {{ display: grid; grid-template-columns: 1.35fr 1fr; gap: 34px; margin-top: 4px; flex: 1; }}
  .highlight {{
    background: var(--panel); border: 1px solid var(--line); border-left: 4px solid var(--accent);
    border-radius: 12px; padding: 24px 24px; align-self: start;
  }}
  .highlight h3 {{ color: var(--accent); margin-bottom: 10px; font-size: 16px;
                   text-transform: uppercase; letter-spacing: .07em; }}
  .highlight p {{ font-size: 16px; line-height: 1.6; }}
  table {{ width: 100%; border-collapse: collapse; margin-top: 8px; }}
  th {{
    text-align: left; font-size: 13px; text-transform: uppercase; letter-spacing: .07em;
    color: var(--muted); padding: 12px 14px; border-bottom: 2px solid var(--line);
  }}
  td {{ font-size: 16px; padding: 13px 14px; border-bottom: 1px solid var(--line); }}
  tr:last-child td {{ border-bottom: none; }}
  .shot {{ flex: 1; display: flex; align-items: flex-start; justify-content: center; margin-top: 10px; }}
  .shot img {{
    max-width: 100%; max-height: 100%; border: 1px solid var(--line);
    border-radius: 10px; box-shadow: 0 10px 30px rgba(0,0,0,.28); object-fit: contain;
  }}
  .notice {{ justify-content: center; }}
  .notice .box {{
    background: var(--panel); border: 1px solid var(--accent); border-radius: 16px; padding: 34px 38px;
  }}
  .notice .box p.main {{ font-size: 22px; line-height: 1.55; color: var(--title); margin-bottom: 20px; }}
  .notice .box p.main b {{ color: var(--accent); }}
  .notice ul li {{ font-size: 16px; color: var(--muted); }}
  .notice ul li::before {{ background: var(--accent); }}
  .steps {{ counter-reset: s; display: flex; flex-direction: column; gap: 14px; margin-top: 6px; }}
  .step {{ display: flex; gap: 16px; align-items: flex-start; font-size: 19px; }}
  .step .k {{
    counter-increment: s; flex: none; width: 30px; height: 30px; border-radius: 9px;
    background: var(--brand); color: #fff; font-weight: 700; font-size: 15px;
    display: grid; place-items: center;
  }}
  .step .k::before {{ content: counter(s); }}
  .note {{
    margin-top: 26px; font-size: 15px; color: var(--muted); border-left: 3px solid var(--accent);
    padding-left: 14px; line-height: 1.55;
  }}
  .bar {{
    position: fixed; left: 0; right: 0; bottom: 0; background: rgba(5,8,15,.94);
    border-top: 1px solid #1b2740; padding: 10px 18px; display: flex; gap: 12px;
    align-items: center; justify-content: center; font-size: 13px; color: #8FA3BF;
    backdrop-filter: blur(6px);
  }}
  .bar button {{
    background: #17263E; color: #E6EDF7; border: 1px solid #22304A; border-radius: 8px;
    padding: 7px 14px; font-size: 13px; cursor: pointer; font-family: inherit;
  }}
  .bar button:hover {{ background: #1F6FEB; border-color: #1F6FEB; }}
  @media print {{
    body {{ background: #fff; padding: 0; }}
    .bar {{ display: none; }}
    .deck {{ gap: 0; width: 100%; }}
    .slide {{
      border: none; border-radius: 0; page-break-after: always; break-after: page;
      width: 100%; aspect-ratio: auto; height: 100vh;
    }}
  }}
  @page {{ size: landscape; margin: 0; }}
</style>
</head>
<body>
<div class="deck">
{slides}
</div>
<div class="bar">
  <button onclick="jump(-1)">&larr; {prev_label}</button>
  <span id="pos">1 / {count}</span>
  <button onclick="jump(1)">{next_label} &rarr;</button>
  <button onclick="window.print()">{print_label}</button>
</div>
<script>
  const slides = [...document.querySelectorAll('.slide')];
  let index = 0;
  function show(i) {{
    index = Math.max(0, Math.min(slides.length - 1, i));
    slides[index].scrollIntoView({{ behavior: 'smooth', block: 'center' }});
    document.getElementById('pos').textContent = (index + 1) + ' / ' + slides.length;
  }}
  function jump(step) {{ show(index + step); }}
  document.addEventListener('keydown', (e) => {{
    if (['ArrowRight', 'PageDown', ' '].includes(e.key)) {{ e.preventDefault(); jump(1); }}
    if (['ArrowLeft', 'PageUp'].includes(e.key)) {{ e.preventDefault(); jump(-1); }}
    if (e.key === 'Home') {{ e.preventDefault(); show(0); }}
    if (e.key === 'End') {{ e.preventDefault(); show(slides.length - 1); }}
  }});
  // Kaydırdıkça göstergeyi güncelle
  const observer = new IntersectionObserver((entries) => {{
    entries.forEach((entry) => {{
      if (entry.isIntersecting) {{
        index = slides.indexOf(entry.target);
        document.getElementById('pos').textContent = (index + 1) + ' / ' + slides.length;
      }}
    }});
  }}, {{ threshold: 0.5 }});
  slides.forEach((s) => observer.observe(s));
</script>
</body>
</html>
"""


def esc(text: str) -> str:
    return html_module.escape(strip_markup(text))


def esc_emphasis(text: str) -> str:
    return "".join(
        f"<b>{html_module.escape(chunk)}</b>" if bold else html_module.escape(chunk)
        for chunk, bold in emphasis_parts(text)
    )


def build_html(language: str, printable: bool) -> str:
    colors = palette(printable)
    meta = META[language]
    blocks: list[str] = []

    for number, slide in enumerate(SLIDES, start=1):
        data = slide[language]
        kind = slide["kind"]
        body = ""

        if kind == "cover":
            chips = "".join(f'<span class="chip">{esc(c)}</span>' for c in data["chips"])
            body = (
                f'<div class="mark">🍽️</div>'
                f'<h1>{esc(data["title"])}</h1>'
                f'<div class="sub">{esc(data["subtitle"])}</div>'
                f'<p class="tag">{esc(data["tagline"])}</p>'
                f'<div class="chips">{chips}</div>'
            )
        elif kind == "bullets":
            items = "".join(f"<li>{esc_emphasis(b)}</li>" for b in data["bullets"])
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["lead"])}</p>'
                f"<ul>{items}</ul>"
            )
        elif kind == "stats":
            cards = "".join(
                f'<div class="stat"><div class="v">{esc(s["value"])}</div>'
                f'<div class="l">{esc(s["label"])}</div>'
                f'<div class="n">{esc(s["note"])}</div></div>'
                for s in data["stats"]
            )
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["lead"])}</p>'
                f'<div class="stats">{cards}</div>'
            )
        elif kind == "grid":
            cards = "".join(
                f'<div class="card"><h3>{esc(i["title"])}</h3><p>{esc(i["text"])}</p></div>'
                for i in data["items"]
            )
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["lead"])}</p>'
                f'<div class="grid">{cards}</div>'
            )
        elif kind == "split":
            items = "".join(f"<li>{esc_emphasis(b)}</li>" for b in data["bullets"])
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<div class="split"><ul>{items}</ul>'
                f'<div class="highlight"><h3>{esc(data["highlight_title"])}</h3>'
                f'<p>{esc(data["highlight_text"])}</p></div></div>'
            )
        elif kind == "table":
            head = "".join(f"<th>{esc(h)}</th>" for h in data["headers"])
            rows = "".join(
                "<tr>" + "".join(f"<td>{esc(c)}</td>" for c in row) + "</tr>"
                for row in data["rows"]
            )
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["lead"])}</p>'
                f"<table><thead><tr>{head}</tr></thead><tbody>{rows}</tbody></table>"
            )
        elif kind == "notice":
            details = "".join(f"<li>{esc(d)}</li>" for d in data["details"])
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<div class="box"><p class="main">{esc_emphasis(data["text"])}</p>'
                f"<ul>{details}</ul></div>"
            )
        elif kind == "screenshot":
            # HTML çıktısı docs/presentation/ altında; ekran görüntüleri
            # sunum/screenshots/ altında durur.
            image_src = f"../../sunum/screenshots/{language}/{slide['image']}"
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["caption"])}</p>'
                f'<div class="shot"><img src="{image_src}" alt="{esc(data["title"])}"></div>'
            )
        elif kind == "closing":
            steps = "".join(
                f'<div class="step"><span class="k"></span><span>{esc(s)}</span></div>'
                for s in data["steps"]
            )
            body = (
                f'<h2>{esc(data["title"])}</h2>'
                f'<p class="lead">{esc(data["lead"])}</p>'
                f'<div class="steps">{steps}</div>'
                f'<p class="note">{esc(data["note"])}</p>'
            )

        css_class = f"slide {kind}" if kind in {"cover", "notice"} else "slide"
        blocks.append(
            f'<section class="{css_class}">{body}'
            f'<div class="foot"><span>{esc(meta["footer"])}</span>'
            f"<span>{number} / {len(SLIDES)}</span></div></section>"
        )

    labels = {
        "tr": ("Önceki", "Sonraki", "Yazdır"),
        "en": ("Previous", "Next", "Print"),
    }[language]

    return HTML_TEMPLATE.format(
        lang=language,
        title=f'{meta["product"]} — {meta["subtitle"]}',
        slides="\n".join(blocks),
        count=len(SLIDES),
        prev_label=labels[0],
        next_label=labels[1],
        print_label=labels[2],
        bg=colors["bg"],
        panel=colors["panel"],
        title_color=colors["title"],
        text=colors["text"],
        muted=colors["muted"],
        brand=colors["brand"],
        accent=colors["accent"],
        line=colors["line"],
        chip_bg=colors["chip_bg"],
    )


# ==================================================================
#  PowerPoint
# ==================================================================
def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def add_textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    return frame


def style_run(run, *, size, color, bold=False, font="Segoe UI"):
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.name = font


def add_rect(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.word_wrap = True
    return shape


def build_pptx(language: str, printable: bool, path: Path) -> None:
    colors = palette(printable)
    meta = META[language]

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W_IN)
    presentation.slide_height = Inches(SLIDE_H_IN)
    blank = presentation.slide_layouts[6]

    margin = Inches(0.72)
    content_width = presentation.slide_width - margin * 2

    for number, slide_data in enumerate(SLIDES, start=1):
        data = slide_data[language]
        kind = slide_data["kind"]
        slide = presentation.slides.add_slide(blank)

        # Zemin
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = rgb(colors["bg"])

        # Alt vurgu şeridi
        strip = add_rect(
            slide,
            0,
            presentation.slide_height - Emu(45720),
            presentation.slide_width,
            Emu(45720),
            colors["brand"],
        )
        strip.adjustments[0] = 0

        top = margin

        if kind == "cover":
            frame = add_textbox(slide, margin, Inches(2.1), content_width, Inches(1.4))
            run = frame.paragraphs[0].add_run()
            run.text = data["title"]
            style_run(run, size=54, color=colors["title"], bold=True)

            frame = add_textbox(slide, margin, Inches(3.35), content_width, Inches(0.7))
            run = frame.paragraphs[0].add_run()
            run.text = data["subtitle"]
            style_run(run, size=26, color=colors["muted"])

            frame = add_textbox(slide, margin, Inches(4.15), Inches(9.2), Inches(1.2))
            run = frame.paragraphs[0].add_run()
            run.text = data["tagline"]
            style_run(run, size=18, color=colors["text"])

            chip_left = margin
            for chip in data["chips"]:
                width = Inches(0.32 + len(chip) * 0.105)
                shape = add_rect(
                    slide,
                    chip_left,
                    Inches(5.55),
                    width,
                    Inches(0.45),
                    colors["chip_bg"],
                    colors["line"],
                )
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = chip
                style_run(run, size=12, color=colors["text"], bold=True)
                chip_left += width + Inches(0.14)

        else:
            frame = add_textbox(slide, margin, top, content_width, Inches(0.8))
            run = frame.paragraphs[0].add_run()
            run.text = data["title"]
            style_run(run, size=32, color=colors["title"], bold=True)
            top += Inches(0.85)

            if data.get("lead"):
                frame = add_textbox(slide, margin, top, Inches(10.5), Inches(0.7))
                run = frame.paragraphs[0].add_run()
                run.text = data["lead"]
                style_run(run, size=15, color=colors["muted"])
                top += Inches(0.78)

        if kind == "bullets":
            frame = add_textbox(slide, margin, top, content_width, Inches(4.2))
            for index, bullet in enumerate(data["bullets"]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.space_after = Pt(11)
                marker = paragraph.add_run()
                marker.text = "▪  "
                style_run(marker, size=16, color=colors["brand"], bold=True)
                run = paragraph.add_run()
                run.text = strip_markup(bullet)
                style_run(run, size=16, color=colors["text"])

        elif kind == "stats":
            gap = Inches(0.24)
            width = (content_width - gap * 3) / 4
            for index, stat in enumerate(data["stats"]):
                left = margin + (width + gap) * index
                card = add_rect(
                    slide, left, top, width, Inches(2.35), colors["panel"], colors["line"]
                )
                frame = card.text_frame
                frame.margin_left = Inches(0.24)
                frame.margin_top = Inches(0.26)
                frame.margin_right = Inches(0.18)

                paragraph = frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = stat["value"]
                style_run(run, size=40, color=colors["brand"], bold=True)

                paragraph = frame.add_paragraph()
                paragraph.space_before = Pt(6)
                run = paragraph.add_run()
                run.text = stat["label"]
                style_run(run, size=14, color=colors["title"], bold=True)

                paragraph = frame.add_paragraph()
                paragraph.space_before = Pt(3)
                run = paragraph.add_run()
                run.text = stat["note"]
                style_run(run, size=11, color=colors["muted"])

        elif kind == "grid":
            columns, gap = 4, Inches(0.2)
            width = (content_width - gap * (columns - 1)) / columns
            height = Inches(1.28)
            for index, item in enumerate(data["items"]):
                row, column = divmod(index, columns)
                left = margin + (width + gap) * column
                card_top = top + (height + Inches(0.16)) * row
                card = add_rect(
                    slide, left, card_top, width, height, colors["panel"], colors["line"]
                )
                frame = card.text_frame
                frame.margin_left = Inches(0.16)
                frame.margin_top = Inches(0.14)
                frame.margin_right = Inches(0.12)

                paragraph = frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = item["title"]
                style_run(run, size=12.5, color=colors["title"], bold=True)

                paragraph = frame.add_paragraph()
                paragraph.space_before = Pt(3)
                run = paragraph.add_run()
                run.text = item["text"]
                style_run(run, size=10, color=colors["muted"])

        elif kind == "split":
            left_width = content_width * 0.56
            right_width = content_width - left_width - Inches(0.4)

            frame = add_textbox(slide, margin, top, left_width, Inches(4.0))
            for index, bullet in enumerate(data["bullets"]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.space_after = Pt(10)
                marker = paragraph.add_run()
                marker.text = "▪  "
                style_run(marker, size=15, color=colors["brand"], bold=True)
                run = paragraph.add_run()
                run.text = strip_markup(bullet)
                style_run(run, size=15, color=colors["text"])

            card = add_rect(
                slide,
                margin + left_width + Inches(0.4),
                top,
                right_width,
                Inches(3.1),
                colors["panel"],
                colors["accent"],
            )
            frame = card.text_frame
            frame.margin_left = Inches(0.26)
            frame.margin_top = Inches(0.24)
            frame.margin_right = Inches(0.22)

            paragraph = frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = upper(data["highlight_title"], language)
            style_run(run, size=12, color=colors["accent"], bold=True)

            paragraph = frame.add_paragraph()
            paragraph.space_before = Pt(8)
            run = paragraph.add_run()
            run.text = data["highlight_text"]
            style_run(run, size=13, color=colors["text"])

        elif kind == "table":
            rows = len(data["rows"]) + 1
            columns = len(data["headers"])
            shape = slide.shapes.add_table(
                rows, columns, margin, top, content_width, Inches(0.5 * rows)
            )
            table = shape.table
            for index, header in enumerate(data["headers"]):
                cell = table.cell(0, index)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(colors["brand"])
                paragraph = cell.text_frame.paragraphs[0]
                style_run(paragraph.runs[0], size=13, color="#FFFFFF", bold=True)

            for row_index, row in enumerate(data["rows"], start=1):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = value
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(colors["panel"])
                    paragraph = cell.text_frame.paragraphs[0]
                    style_run(paragraph.runs[0], size=12, color=colors["text"])

        elif kind == "screenshot":
            frame = add_textbox(slide, margin, top, Inches(11.0), Inches(0.55))
            run = frame.paragraphs[0].add_run()
            run.text = data["caption"]
            style_run(run, size=14, color=colors["muted"])
            top += Inches(0.6)

            image_path = SCREENSHOT_DIR / language / slide_data["image"]
            if image_path.is_file():
                available = Inches(SLIDE_H_IN) - top - Inches(0.35)
                image_width = Emu(int(available * 16 / 9))
                if image_width > content_width:
                    image_width = content_width
                    available = Emu(int(image_width * 9 / 16))
                left = (presentation.slide_width - image_width) // 2
                picture = slide.shapes.add_picture(
                    str(image_path), left, top, width=image_width, height=available
                )
                picture.line.color.rgb = rgb(colors["line"])
                picture.line.width = Pt(0.75)
            else:
                print(f"  [!] Görsel yok, slayt boş kalacak: {image_path}")

        elif kind == "notice":
            card = add_rect(
                slide,
                margin,
                Inches(1.9),
                content_width,
                Inches(3.9),
                colors["panel"],
                colors["accent"],
            )
            frame = card.text_frame
            frame.margin_left = Inches(0.4)
            frame.margin_top = Inches(0.36)
            frame.margin_right = Inches(0.4)

            paragraph = frame.paragraphs[0]
            for chunk, bold in emphasis_parts(data["text"]):
                run = paragraph.add_run()
                run.text = chunk
                style_run(
                    run,
                    size=17,
                    color=colors["accent"] if bold else colors["title"],
                    bold=bold,
                )

            for detail in data["details"]:
                paragraph = frame.add_paragraph()
                paragraph.space_before = Pt(12)
                marker = paragraph.add_run()
                marker.text = "▪  "
                style_run(marker, size=13, color=colors["accent"], bold=True)
                run = paragraph.add_run()
                run.text = detail
                style_run(run, size=13, color=colors["muted"])

        elif kind == "closing":
            frame = add_textbox(slide, margin, top, content_width, Inches(3.0))
            for index, step in enumerate(data["steps"]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.space_after = Pt(12)
                marker = paragraph.add_run()
                marker.text = f"{index + 1}.  "
                style_run(marker, size=17, color=colors["brand"], bold=True)
                run = paragraph.add_run()
                run.text = step
                style_run(run, size=17, color=colors["text"])

            frame = add_textbox(slide, margin, Inches(5.9), content_width, Inches(0.8))
            run = frame.paragraphs[0].add_run()
            run.text = data["note"]
            style_run(run, size=13, color=colors["muted"])

        # Altbilgi
        frame = add_textbox(
            slide, margin, presentation.slide_height - Inches(0.56), content_width, Inches(0.3)
        )
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = meta["footer"]
        style_run(run, size=9, color=colors["muted"])

        frame = add_textbox(
            slide,
            presentation.slide_width - margin - Inches(1.4),
            presentation.slide_height - Inches(0.56),
            Inches(1.4),
            Inches(0.3),
        )
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        run = paragraph.add_run()
        run.text = f"{number} / {len(SLIDES)}"
        style_run(run, size=9, color=colors["muted"])

    presentation.save(str(path))


# ==================================================================
#  PDF
# ==================================================================
def build_pdf(language: str, printable: bool, path: Path, fonts: tuple[str, str]) -> None:
    regular, bold = fonts
    colors = palette(printable)
    meta = META[language]
    width, height = PAGE_SIZE
    margin = 52

    canvas = pdf_canvas.Canvas(str(path), pagesize=PAGE_SIZE)
    # PDF üstverisi de yayımlanan bir yüzeydir: kullanıcı adı, makine adı ya
    # da yerel dosya yolu sızdırmamalıdır. Yalnızca ürün bilgisi yazılır.
    canvas.setTitle(f'{meta["product"]} — {meta["subtitle"]}')
    canvas.setAuthor(meta["footer"])
    canvas.setSubject(meta["tagline"])
    canvas.setCreator(meta["footer"])
    canvas.setKeywords("")

    def fill(color: str) -> None:
        canvas.setFillColor(HexColor(color))

    def text_lines(lines, x, y, font, size, color, leading):
        fill(color)
        canvas.setFont(font, size)
        for line in lines:
            canvas.drawString(x, y, line)
            y -= leading
        return y

    def rounded(x, y, w, h, fill_color, stroke=None):
        fill(fill_color)
        if stroke:
            canvas.setStrokeColor(HexColor(stroke))
            canvas.setLineWidth(0.7)
            canvas.roundRect(x, y, w, h, 8, stroke=1, fill=1)
        else:
            canvas.roundRect(x, y, w, h, 8, stroke=0, fill=1)

    for number, slide_data in enumerate(SLIDES, start=1):
        data = slide_data[language]
        kind = slide_data["kind"]

        fill(colors["bg"])
        canvas.rect(0, 0, width, height, stroke=0, fill=1)
        fill(colors["brand"])
        canvas.rect(0, 0, width, 5, stroke=0, fill=1)

        y = height - margin - 20
        body_width = width - margin * 2

        if kind == "cover":
            fill(colors["brand"])
            canvas.roundRect(margin, height - 190, 58, 58, 14, stroke=0, fill=1)

            y = height - 250
            fill(colors["title"])
            canvas.setFont(bold, 44)
            canvas.drawString(margin, y, data["title"])

            y -= 44
            fill(colors["muted"])
            canvas.setFont(regular, 22)
            canvas.drawString(margin, y, data["subtitle"])

            y -= 46
            for line in wrap_text(data["tagline"], regular, 15, body_width * 0.72):
                fill(colors["text"])
                canvas.setFont(regular, 15)
                canvas.drawString(margin, y, line)
                y -= 22

            y -= 18
            chip_x = margin
            for chip in data["chips"]:
                chip_w = pdfmetrics.stringWidth(chip, bold, 10) + 26
                rounded(chip_x, y - 6, chip_w, 24, colors["chip_bg"], colors["line"])
                fill(colors["text"])
                canvas.setFont(bold, 10)
                canvas.drawString(chip_x + 13, y + 1, chip)
                chip_x += chip_w + 9

        else:
            fill(colors["title"])
            canvas.setFont(bold, 26)
            canvas.drawString(margin, y, data["title"])
            y -= 30

            if data.get("lead"):
                for line in wrap_text(data["lead"], regular, 12.5, body_width * 0.78):
                    fill(colors["muted"])
                    canvas.setFont(regular, 12.5)
                    canvas.drawString(margin, y, line)
                    y -= 17
            y -= 12

        if kind == "bullets":
            for bullet in data["bullets"]:
                lines = wrap_text(strip_markup(bullet), regular, 13.5, body_width - 22)
                fill(colors["brand"])
                canvas.rect(margin + 1, y + 3, 6, 6, stroke=0, fill=1)
                y = text_lines(lines, margin + 20, y, regular, 13.5, colors["text"], 18)
                y -= 8

        elif kind == "stats":
            gap = 16
            card_w = (body_width - gap * 3) / 4
            card_h = 150
            card_y = y - card_h
            for index, stat in enumerate(data["stats"]):
                x = margin + (card_w + gap) * index
                rounded(x, card_y, card_w, card_h, colors["panel"], colors["line"])
                fill(colors["brand"])
                canvas.setFont(bold, 34)
                canvas.drawString(x + 18, card_y + card_h - 48, stat["value"])
                fill(colors["title"])
                canvas.setFont(bold, 12)
                canvas.drawString(x + 18, card_y + card_h - 74, stat["label"])
                note_y = card_y + card_h - 94
                for line in wrap_text(stat["note"], regular, 9.5, card_w - 36):
                    fill(colors["muted"])
                    canvas.setFont(regular, 9.5)
                    canvas.drawString(x + 18, note_y, line)
                    note_y -= 12

        elif kind == "grid":
            columns, gap = 4, 14
            card_w = (body_width - gap * (columns - 1)) / columns
            card_h = 84
            for index, item in enumerate(data["items"]):
                row, column = divmod(index, columns)
                x = margin + (card_w + gap) * column
                card_y = y - card_h - (card_h + 12) * row
                rounded(x, card_y, card_w, card_h, colors["panel"], colors["line"])
                fill(colors["title"])
                canvas.setFont(bold, 10.5)
                canvas.drawString(x + 12, card_y + card_h - 22, item["title"][:34])
                note_y = card_y + card_h - 38
                for line in wrap_text(item["text"], regular, 8.5, card_w - 24)[:4]:
                    fill(colors["muted"])
                    canvas.setFont(regular, 8.5)
                    canvas.drawString(x + 12, note_y, line)
                    note_y -= 11

        elif kind == "split":
            left_w = body_width * 0.55
            bullet_y = y
            for bullet in data["bullets"]:
                lines = wrap_text(strip_markup(bullet), regular, 12.5, left_w - 22)
                fill(colors["brand"])
                canvas.rect(margin + 1, bullet_y + 3, 6, 6, stroke=0, fill=1)
                bullet_y = text_lines(
                    lines, margin + 20, bullet_y, regular, 12.5, colors["text"], 17
                )
                bullet_y -= 7

            box_x = margin + left_w + 28
            box_w = body_width - left_w - 28
            wrapped = wrap_text(data["highlight_text"], regular, 11.5, box_w - 34)
            box_h = 56 + len(wrapped) * 16
            box_y = y - box_h + 18
            rounded(box_x, box_y, box_w, box_h, colors["panel"], colors["line"])
            fill(colors["accent"])
            canvas.rect(box_x, box_y, 3.5, box_h, stroke=0, fill=1)

            fill(colors["accent"])
            canvas.setFont(bold, 10.5)
            canvas.drawString(
                box_x + 18, box_y + box_h - 26, upper(data["highlight_title"], language)
            )
            note_y = box_y + box_h - 48
            for line in wrapped:
                fill(colors["text"])
                canvas.setFont(regular, 11.5)
                canvas.drawString(box_x + 18, note_y, line)
                note_y -= 16

        elif kind == "table":
            columns = len(data["headers"])
            column_w = body_width / columns
            header_h = 28
            row_h = 26

            fill(colors["brand"])
            canvas.rect(margin, y - header_h, body_width, header_h, stroke=0, fill=1)
            for index, header in enumerate(data["headers"]):
                fill("#FFFFFF")
                canvas.setFont(bold, 10.5)
                canvas.drawString(margin + column_w * index + 12, y - 19, header)

            row_y = y - header_h
            for row_index, row in enumerate(data["rows"]):
                row_y -= row_h
                if row_index % 2 == 0:
                    fill(colors["panel"])
                    canvas.rect(margin, row_y, body_width, row_h, stroke=0, fill=1)
                for column_index, value in enumerate(row):
                    fill(colors["text"])
                    canvas.setFont(regular, 10.5)
                    canvas.drawString(margin + column_w * column_index + 12, row_y + 9, value[:44])

        elif kind == "screenshot":
            for line in wrap_text(data["caption"], regular, 12.5, body_width * 0.9):
                fill(colors["muted"])
                canvas.setFont(regular, 12.5)
                canvas.drawString(margin, y, line)
                y -= 17
            y -= 8

            image_path = SCREENSHOT_DIR / language / slide_data["image"]
            if image_path.is_file():
                bottom_limit = 26
                image_h = y - bottom_limit
                image_w = image_h * 16 / 9
                if image_w > body_width:
                    image_w = body_width
                    image_h = image_w * 9 / 16
                x = (width - image_w) / 2
                canvas.drawImage(
                    str(image_path),
                    x,
                    y - image_h,
                    width=image_w,
                    height=image_h,
                    preserveAspectRatio=True,
                    mask="auto",
                )
                canvas.setStrokeColor(HexColor(colors["line"]))
                canvas.setLineWidth(0.8)
                canvas.rect(x, y - image_h, image_w, image_h, stroke=1, fill=0)
            else:
                print(f"  [!] Görsel yok: {image_path}")

        elif kind == "notice":
            wrapped = wrap_text(strip_markup(data["text"]), bold, 15, body_width - 70)
            detail_lines = [
                wrap_text(detail, regular, 11.5, body_width - 90) for detail in data["details"]
            ]
            box_h = 60 + len(wrapped) * 22 + sum(len(d) for d in detail_lines) * 15 + 30
            box_y = y - box_h + 20
            rounded(margin, box_y, body_width, box_h, colors["panel"], colors["accent"])

            note_y = box_y + box_h - 36
            for line in wrapped:
                fill(colors["title"])
                canvas.setFont(bold, 15)
                canvas.drawString(margin + 30, note_y, line)
                note_y -= 22

            note_y -= 10
            for lines in detail_lines:
                fill(colors["accent"])
                canvas.rect(margin + 32, note_y + 3, 5, 5, stroke=0, fill=1)
                for line in lines:
                    fill(colors["muted"])
                    canvas.setFont(regular, 11.5)
                    canvas.drawString(margin + 48, note_y, line)
                    note_y -= 15
                note_y -= 6

        elif kind == "closing":
            for index, step in enumerate(data["steps"], start=1):
                fill(colors["brand"])
                canvas.roundRect(margin, y - 5, 22, 22, 6, stroke=0, fill=1)
                fill("#FFFFFF")
                canvas.setFont(bold, 11)
                canvas.drawCentredString(margin + 11, y + 1, str(index))
                lines = wrap_text(step, regular, 13.5, body_width - 50)
                y = text_lines(lines, margin + 34, y, regular, 13.5, colors["text"], 18)
                y -= 12

            y -= 10
            fill(colors["accent"])
            canvas.rect(margin, y - 24, 3, 34, stroke=0, fill=1)
            note_y = y
            for line in wrap_text(data["note"], regular, 11, body_width - 40):
                fill(colors["muted"])
                canvas.setFont(regular, 11)
                canvas.drawString(margin + 14, note_y, line)
                note_y -= 15

        # Altbilgi
        fill(colors["muted"])
        canvas.setFont(regular, 8.5)
        canvas.drawString(margin, 22, meta["footer"])
        canvas.drawRightString(width - margin, 22, f"{number} / {len(SLIDES)}")

        canvas.showPage()

    canvas.save()


# ==================================================================
#  Ana akış
# ==================================================================
def main() -> int:
    only = None
    if "--only" in sys.argv:
        only = sys.argv[sys.argv.index("--only") + 1]

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    fonts = register_pdf_fonts()

    print()
    print("  Akıllı Restaurant — tanıtım sunumu")
    print("  " + "-" * 46)
    print(f"  Slayt sayısı : {len(SLIDES)}")
    print(f"  Yazı tipi    : {fonts[0]}")
    print(f"  Çıktı        : {OUTPUT_DIR}")
    print()

    produced: list[Path] = []
    for language in ("tr", "en"):
        # Sayılar her dil için yeniden ölçülür (binlik ayracı da dile bağlı).
        use_language(language)
        for variant, printable in (("screen", False), ("print", True)):
            stem = FILENAMES[language][variant]

            if variant == "screen" and only in (None, "html"):
                path = OUTPUT_DIR / f"{stem}.html"
                path.write_text(build_html(language, printable), encoding="utf-8")
                produced.append(path)

            if only in (None, "pptx"):
                path = OUTPUT_DIR / f"{stem}.pptx"
                build_pptx(language, printable, path)
                produced.append(path)

            if only in (None, "pdf"):
                path = OUTPUT_DIR / f"{stem}.pdf"
                build_pdf(language, printable, path, fonts)
                produced.append(path)

    for path in produced:
        size_kb = round(path.stat().st_size / 1024)
        print(f"  [OK] {path.name:<44} {size_kb:>6} KB")

    print()
    print(f"  Toplam {len(produced)} dosya üretildi.")
    print()
    return 0


if __name__ == "__main__":
    sys.exit(main())
